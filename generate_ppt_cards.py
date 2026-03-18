#!/usr/bin/env python3
"""
Generate PPT with card-based layout design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(230, 240, 250)
GRAY = RGBColor(245, 245, 245)
DARK_GRAY = RGBColor(80, 80, 80)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Light gray background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRAY
    bg.line.fill.background()
    
    # Blue bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_card_slide_2col(prs, title, left_title, left_items, right_title, right_items):
    """Two column layout with cards"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Left card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = LIGHT_BLUE
    left_card.line.color.rgb = BLUE
    left_card.line.width = Pt(2)
    
    # Left title
    lt_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5))
    tf = lt_box.text_frame
    tf.text = left_title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Left items
    li_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.8))
    tf = li_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right card
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(5.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_BLUE
    right_card.line.color.rgb = BLUE
    right_card.line.width = Pt(2)
    
    # Right title
    rt_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.6), Inches(0.5))
    tf = rt_box.text_frame
    tf.text = right_title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Right items
    ri_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.6), Inches(4.8))
    tf = ri_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

def add_card_slide_3col(prs, title, col1_title, col1_items, col2_title, col2_items, col3_title, col3_items):
    """Three column layout with cards"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    col_width = 3.9
    gap = 0.25
    
    # Column 1
    x1 = 0.5
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x1), Inches(1.3), Inches(col_width), Inches(5.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(255, 240, 240)
    card1.line.color.rgb = RGBColor(180, 80, 80)
    card1.line.width = Pt(2)
    
    t1_box = slide.shapes.add_textbox(Inches(x1+0.15), Inches(1.5), Inches(col_width-0.3), Inches(0.5))
    tf = t1_box.text_frame
    tf.text = col1_title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 80, 80)
    
    i1_box = slide.shapes.add_textbox(Inches(x1+0.15), Inches(2.1), Inches(col_width-0.3), Inches(4.8))
    tf = i1_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(col1_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Column 2
    x2 = x1 + col_width + gap
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x2), Inches(1.3), Inches(col_width), Inches(5.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(240, 255, 240)
    card2.line.color.rgb = RGBColor(80, 150, 80)
    card2.line.width = Pt(2)
    
    t2_box = slide.shapes.add_textbox(Inches(x2+0.15), Inches(1.5), Inches(col_width-0.3), Inches(0.5))
    tf = t2_box.text_frame
    tf.text = col2_title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 150, 80)
    
    i2_box = slide.shapes.add_textbox(Inches(x2+0.15), Inches(2.1), Inches(col_width-0.3), Inches(4.8))
    tf = i2_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(col2_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Column 3
    x3 = x2 + col_width + gap
    card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x3), Inches(1.3), Inches(col_width), Inches(5.8))
    card3.fill.solid()
    card3.fill.fore_color.rgb = RGBColor(240, 240, 255)
    card3.line.color.rgb = RGBColor(80, 80, 180)
    card3.line.width = Pt(2)
    
    t3_box = slide.shapes.add_textbox(Inches(x3+0.15), Inches(1.5), Inches(col_width-0.3), Inches(0.5))
    tf = t3_box.text_frame
    tf.text = col3_title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 80, 180)
    
    i3_box = slide.shapes.add_textbox(Inches(x3+0.15), Inches(2.1), Inches(col_width-0.3), Inches(4.8))
    tf = i3_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(col3_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    return slide

def add_highlight_slide(prs, title, highlights):
    """Slide with highlighted key numbers"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Create highlight boxes
    box_width = 3.8
    box_height = 2.2
    gap = 0.4
    
    positions = [
        (0.5, 1.3), (4.7, 1.3), (8.9, 1.3),
        (0.5, 3.9), (4.7, 3.9), (8.9, 3.9)
    ]
    
    for i, (number, label) in enumerate(highlights):
        if i >= 6:
            break
        x, y = positions[i]
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = BLUE
        box.line.width = Pt(2)
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(y+0.3), Inches(box_width), Inches(1))
        tf = num_box.text_frame
        tf.text = number
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y+1.3), Inches(box_width), Inches(0.8))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ===== GENERATE SLIDES =====

# Slide 1: Title
add_title_slide(prs, "皇庭国际中心", "数字化营销服务建议书")

# Slide 2: Project Overview - 2 columns
add_card_slide_2col(prs, "项目概况",
    "项目基本信息",
    [
        "总建筑面积：12万㎡",
        "写字楼：8万㎡ | 商业：2万㎡ | 停车：2万㎡",
        "地上38层，地下3层",
        "预计2024年Q3交付",
        "5A甲级智能商务办公",
        "LEED金级预认证"
    ],
    "核心优势",
    [
        "建筑设计：国际知名事务所操刀",
        "单层1500-2000㎡，层高4.2米",
        "全玻璃幕墙+LOW-E节能设计",
        "智能化：人脸识别/智能停车/5G/楼宇自控",
        "双地铁交汇，距地铁站仅200米",
        "第一太平戴维斯24小时管家服务"
    ]
)

# Slide 3: Target Customers - 3 columns
add_card_slide_3col(prs, "目标客群定位",
    "金融机构 35%",
    [
        "银行分行",
        "证券公司营业部",
        "保险公司区域中心",
        "基金公司",
        "租金承受能力强",
        "注重楼宇品质"
    ],
    "科技企业 25%",
    [
        "互联网公司",
        "软件/AI企业",
        "大数据企业",
        "员工年轻化",
        "注重智能化",
        "追求办公舒适度"
    ],
    "专业服务 20%",
    [
        "律师事务所",
        "会计师事务所",
        "咨询公司",
        "广告公司",
        "商务接待需求高",
        "重视周边配套"
    ]
)

# Slide 4: Market Analysis - Highlight numbers
add_highlight_slide(prs, "市场环境分析",
    [
        ("25%", "GDP贡献占比"),
        ("200万㎡", "现有甲级写字楼"),
        ("5000+", "入驻企业数量"),
        ("8%", "区域空置率"),
        ("30%", "企业升级需求"),
        ("12个月", "无新增供应窗口期")
    ]
)

# Slide 5: Competition - 2 columns
add_card_slide_2col(prs, "竞争项目对比",
    "竞争项目A（东侧800米）",
    [
        "2019年入市，10万㎡",
        "出租率85%",
        "日租金9-11元/㎡",
        "优势：品牌知名度高",
        "劣势：楼龄5年，设施老化",
        "劣势：智能化程度低"
    ],
    "皇庭国际中心优势",
    [
        "2024年全新交付",
        "距地铁仅200米",
        "日租金10-12元/㎡",
        "双地铁交汇上盖",
        "第一太平戴维斯物业",
        "智能化程度最高"
    ]
)

# Slide 6: Section
add_section_slide(prs, "数字化营销策略", "全渠道数字化营销体系")

# Slide 7: Digital Strategy - 3 columns
add_card_slide_3col(prs, "数字化营销矩阵",
    "官网小程序",
    [
        "响应式官网设计",
        "VR全景看房系统",
        "实时房源查询",
        "在线预约看房",
        "ChatGPT智能客服",
        "用户行为数据分析"
    ],
    "短视频直播",
    [
        "抖音+视频号+小红书",
        "每周发布3-5条内容",
        "项目宣传片",
        "入驻企业访谈",
        "每周1-2场直播看房",
        "直播专属优惠"
    ],
    "精准投放",
    [
        "微信朋友圈广告",
        "百度搜索+信息流",
        "LinkedIn国际投放",
        "月均预算15-20万",
        "定向28-50岁决策者",
        "金融/科技/专业服务"
    ]
)

# Slide 8: Savills - 2 columns
add_card_slide_2col(prs, "第一太平戴维斯渠道策略",
    "全球资源优势",
    [
        "160年历史，70+国家",
        "中国15个城市分公司",
        "3000+专业员工",
        "服务500+世界500强",
        "全球客户网络",
        "国际企业选址需求"
    ],
    "渠道合作政策",
    [
        "3公里独家代理权",
        "首年租金2个月佣金",
        "业绩达标额外0.5个月",
        "世界500强额外奖励5万",
        "专属服务团队5人",
        "每周驻场2天"
    ]
)

# Save
prs.save('/Users/bob/.openclaw/workspace/皇庭国际中心_图块排版版.pptx')
print("Done! Card-based layout PPT generated.")
