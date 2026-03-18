#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False
import io
import os

output_dir = "/Users/bob/.openclaw/workspace"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

PRIMARY_COLOR = RgbColor(0x1a, 0x5f, 0x9e)
SECONDARY_COLOR = RgbColor(0x2e, 0x8b, 0x57)
ACCENT_COLOR = RgbColor(0xe6, 0x8a, 0x00)
LIGHT_BG = RgbColor(0xf5, 0xf5, 0xf5)
DARK_TEXT = RgbColor(0x33, 0x33, 0x33)
WHITE = RgbColor(0xff, 0xff, 0xff)
GRAY_TEXT = RgbColor(0x66, 0x66, 0x66)

def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape

def add_body_text(slide, text, left, top, width, height, font_size=18, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = align
    return shape

def add_bullet_points(slide, items, left, top, width, height, font_size=16):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
    return shape

def create_chart_image(data, chart_type='bar', title='', colors=None, figsize=(5.5, 3.8)):
    fig, ax = plt.subplots(figsize=figsize, dpi=150)
    if colors is None:
        colors = ['#1a5f9e', '#2e8b57', '#e68a00', '#c44e52', '#8c564b']
    
    if chart_type == 'bar':
        categories = list(data.keys())
        values = list(data.values())
        bars = ax.bar(categories, values, color=colors[:len(categories)], edgecolor='white')
        ax.set_ylabel('数值', fontsize=10)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height, f'{height:.0f}', ha='center', va='bottom', fontsize=9, fontweight='bold')
    elif chart_type == 'line':
        years = list(range(2021, 2026))
        for i, (label, values) in enumerate(data.items()):
            ax.plot(years[:len(values)], values, marker='o', linewidth=2.5, label=label, color=colors[i], markersize=6)
        ax.set_ylabel('数值', fontsize=10)
        ax.legend(loc='best', fontsize=9)
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
    elif chart_type == 'pie':
        sizes = list(data.values())
        labels = [f"{k}\n{v}%" for k, v in data.items()]
        ax.pie(sizes, labels=labels, colors=colors[:len(sizes)], startangle=90, textprops={'fontsize': 9})
        ax.axis('equal')
    elif chart_type == 'horizontal_bar':
        categories = list(data.keys())
        values = list(data.values())
        bars = ax.barh(categories, values, color=colors[:len(categories)], edgecolor='white')
        ax.set_xlabel('数值', fontsize=10)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for i, bar in enumerate(bars):
            width = bar.get_width()
            ax.text(width, bar.get_y() + bar.get_height()/2., f'{width:.0f}', ha='left', va='center', fontsize=9, fontweight='bold')
    
    ax.set_title(title, fontsize=12, fontweight='bold', pad=12, color='#333333')
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', facecolor='white')
    img_stream.seek(0)
    plt.close()
    return img_stream

def add_background_shape(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_decorative_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_metric_card(slide, title, value, change, left, top, width, height, color=PRIMARY_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    add_title_shape(slide, title, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.35), font_size=12, bold=False, color=GRAY_TEXT)
    add_title_shape(slide, value, left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), Inches(0.5), font_size=22, color=color)
    change_color = SECONDARY_COLOR if "+" in str(change) else RgbColor(0xc4, 0x4e, 0x52)
    add_body_text(slide, change, left + Inches(0.1), top + Inches(0.85), width - Inches(0.2), Inches(0.3), font_size=10, color=change_color)

# 第1页：封面
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, PRIMARY_COLOR)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT_COLOR)
add_decorative_bar(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.08), ACCENT_COLOR)
add_title_shape(slide, "成都写字楼市场", Inches(0.8), Inches(3.0), Inches(11), Inches(0.9), font_size=56, color=WHITE)
add_title_shape(slide, "深度分析报告", Inches(0.8), Inches(3.9), Inches(11), Inches(0.8), font_size=48, color=WHITE)
add_title_shape(slide, "2024-2025年度市场洞察与投资指南", Inches(0.8), Inches(5.0), Inches(11), Inches(0.6), font_size=22, bold=False, color=RgbColor(0xcc, 0xcc, 0xcc))
add_title_shape(slide, "2025年3月", Inches(0.8), Inches(6.3), Inches(4), Inches(0.4), font_size=16, bold=False, color=RgbColor(0xaa, 0xaa, 0xaa))

# 第2页：目录
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "报告目录", Inches(0.6), Inches(0.22), Inches(4), Inches(0.7), font_size=34, color=WHITE)

sections = [("01", "宏观环境", "城市经济、产业背景、政策分析"), ("02", "市场供给", "存量规模、新增供应、项目盘点"), ("03", "市场需求", "行业结构、租户画像、吸纳分析"), ("04", "市场表现", "租金走势、空置率、分区对比"), ("05", "竞争格局", "重点楼宇、运营商分析"), ("06", "趋势预测", "2025-2027市场展望"), ("07", "投资建议", "策略建议与风险提示")]
for i, (num, title, desc) in enumerate(sections):
    row, col = i // 2, i % 2
    left, top = Inches(0.7 + col * 6.3), Inches(1.4 + row * 1.55)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_COLOR if i < 4 else SECONDARY_COLOR
    circle.line.fill.background()
    add_title_shape(slide, num, left, top + Inches(0.12), Inches(0.6), Inches(0.5), font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_shape(slide, title, left + Inches(0.8), top, Inches(2.5), Inches(0.5), font_size=22, color=DARK_TEXT)
    add_body_text(slide, desc, left + Inches(0.8), top + Inches(0.45), Inches(4.8), Inches(0.5), font_size=13, color=RgbColor(0x88, 0x88, 0x88))

# 第3页：城市经济
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "城市经济与产业背景", Inches(0.6), Inches(0.22), Inches(8), Inches(0.7), font_size=32, color=WHITE)
add_title_shape(slide, "2024年核心经济指标", Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.45), font_size=16, color=PRIMARY_COLOR)

metrics = [("GDP总量", "2.35万亿元", "同比+6.0%"), ("常住人口", "2140万人", "净流入+25万"), ("第三产业占比", "65.2%", "同比+0.8pp"), ("世界500强落户", "312家", "新增18家")]
for i, (title, value, change) in enumerate(metrics):
    add_metric_card(slide, title, value, change, Inches(0.5 + (i % 2) * 2.8), Inches(1.85 + (i // 2) * 1.15), Inches(2.6), Inches(1.0), PRIMARY_COLOR)

add_title_shape(slide, "产业结构分布", Inches(7.2), Inches(1.35), Inches(5), Inches(0.45), font_size=16, color=PRIMARY_COLOR)
industry_data = {'电子信息': 28, '装备制造': 22, '金融服务': 15, '生物医药': 12, '其他服务业': 23}
chart_img = create_chart_image(industry_data, 'pie', '', ['#1a5f9e', '#2e8b57', '#e68a00', '#c44e52', '#8c564b'])
slide.shapes.add_picture(chart_img, Inches(6.8), Inches(1.8), width=Inches(6))

# 第4页：市场存量
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "市场存量规模分析", Inches(0.6), Inches(0.22), Inches(8), Inches(0.7), font_size=32, color=WHITE)
add_title_shape(slide, "截至2024年底市场存量", Inches(0.5), Inches(1.35), Inches(6), Inches(0.4), font_size=15, color=PRIMARY_COLOR)

stock_metrics = [("甲级写字楼", "520万㎡", "占比58%"), ("乙级写字楼", "380万㎡", "占比42%"), ("总存量", "900万㎡", "同比+4.2%")]
for i, (title, value, pct) in enumerate(stock_metrics):
    add_metric_card(slide, title, value, pct, Inches(0.5 + i * 2.9), Inches(1.8), Inches(2.7), Inches(0.95), PRIMARY_COLOR if i < 2 else ACCENT_COLOR)

add_title_shape(slide, "存量区域分布（万㎡）", Inches(0.5), Inches(2.95), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
region_data = {'金融城': 120, '天府新区': 110, '大源': 95, '天府广场': 85, '其他': 110}
chart_img = create_chart_image(region_data, 'horizontal_bar', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(3.35), width=Inches(6))

add_title_shape(slide, "历年新增供应趋势（万㎡）", Inches(7), Inches(2.95), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
supply_data = {'新增供应': [58, 52, 45, 38, 35]}
chart_img = create_chart_image(supply_data, 'line', '')
slide.shapes.add_picture(chart_img, Inches(6.8), Inches(3.35), width=Inches(6.2))

# 第5页：新增项目
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "2024-2025年新增供应项目", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)

rows, cols = 5, 6
table = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.35), Inches(12.5), Inches(3.8)).table
headers = ['项目名称', '区域', '体量(万㎡)', '定位', '预租率', '主要租户类型']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    p = cell.text_frame.paragraphs[0]
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

projects = [['西部金融中心', '金融城', '18', '超甲级', '65%', '银行、证券'], ['招商大魔方二期', '大源', '15', '甲级', '55%', '科技、专业服务'], ['天府国金中心北区', '天府新区', '22', '区域总部', '40%', '央企、国企'], ['天府新区CBD项目群', '天府新区', '35', '甲级', '2025入市', '待定']]
for row_idx, row_data in enumerate(projects, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(0xf0, 0xf5, 0xfa)

add_body_text(slide, "注：2025年预计新增供应35万㎡，主要集中在天府新区秦皇寺CBD区域", Inches(0.4), Inches(5.3), Inches(12), Inches(0.4), font_size=12, color=GRAY_TEXT)

# 第6页：需求结构
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "租赁需求结构深度分析", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)
add_title_shape(slide, "需求行业占比", Inches(0.5), Inches(1.35), Inches(5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)

demand_data = {'金融业': 35, '科技互联网': 28, '专业服务': 20, '房地产建筑': 10, '其他': 7}
chart_img = create_chart_image(demand_data, 'pie', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.75), width=Inches(5.8))

add_title_shape(slide, "历年净吸纳量（万㎡）", Inches(6.8), Inches(1.35), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
absorb_data = {'净吸纳量': [42, 38, 35, 32, 30]}
chart_img = create_chart_image(absorb_data, 'line', '')
slide.shapes.add_picture(chart_img, Inches(6.5), Inches(1.75), width=Inches(6.5))

add_title_shape(slide, "重点行业特征", Inches(0.5), Inches(5.0), Inches(12), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
industry_info = ["金融业(35%)：银行、保险、证券、金融科技，偏好金融城、天府广场", "科技互联网(28%)：互联网大厂、本土科技企业，偏好大源、天府新区", "专业服务(20%)：律所、会计所、咨询公司，偏好核心商务区"]
add_bullet_points(slide, industry_info, Inches(0.5), Inches(5.4), Inches(12), Inches(1.8), font_size=13)

# 第7页：租金与空置率
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "租金水平与空置率分析", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)
add_title_shape(slide, "2024年Q4分区平均租金（元/㎡/月）", Inches(0.5), Inches(1.35), Inches(6), Inches(0.4), font_size=15, color=PRIMARY_COLOR)

rent_data = {'金融城': 125, '天府广场': 110, '大源': 85, '天府新区': 75, '其他': 70}
chart_img = create_chart_image(rent_data, 'bar', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.8), width=Inches(6.2))

add_title_shape(slide, "分区空置率对比", Inches(7), Inches(1.35), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
vacancy_data = {'天府广场': 16, '金融城': 18, '大源': 26, '其他': 30, '天府新区': 32}
chart_img = create_chart_image(vacancy_data, 'bar', '', ['#2e8b57', '#5a9', '#e6a23c', '#e68a00', '#c44e52'])
slide.shapes.add_picture(chart_img, Inches(6.8), Inches(1.8), width=Inches(6.2))

add_decorative_bar(slide, Inches(0), Inches(6.0), Inches(13.333), Inches(1.5), LIGHT_BG)
key_metrics = [("全市平均租金", "95元/㎡/月", "同比-4%"), ("甲级空置率", "23.5%", "同比+1.2pp"), ("乙级空置率", "28.0%", "同比+0.8pp"), ("平均去化周期", "20个月", "同比+3个月")]
for i, (title, value, change) in enumerate(key_metrics):
    left = Inches(0.5 + i * 3.2)
    add_title_shape(slide, title, left, Inches(6.15), Inches(2.8), Inches(0.35), font_size=13, bold=False, color=GRAY_TEXT)
    add_title_shape(slide, value, left, Inches(6.45), Inches(2.8), Inches(0.5), font_size=24, color=PRIMARY_COLOR)
    add_body_text(slide, change, left, Inches(6.9), Inches(2.5), Inches(0.35), font_size=11, color=RgbColor(0xc4, 0x4e, 0x52))

# 第8页：区域对比
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "四大商务区深度对比", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)

rows, cols = 9, 5
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8)).table
headers = ['对比维度', '天府广场CBD', '金融城', '大源商务区', '天府新区']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    p = cell.text_frame.paragraphs[0]
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

comparison_data = [['区域定位', '传统CBD', '金融总部商务区', '科技商务新区', '未来城市新中心'], ['发展成熟度', '成熟', '成熟', '较成熟', '发展中'], ['甲级存量(万㎡)', '85', '120', '95', '110'], ['平均租金(元/㎡/月)', '110', '125', '85', '75'], ['空置率', '16%', '18%', '26%', '32%'], ['去化周期', '12个月', '14个月', '22个月', '30个月'], ['主力租户', '金融、专业服务', '金融机构', '科技企业', '央企、国企'], ['核心优势', '配套完善', '产业集聚', '成本适中', '政策优惠']]
for row_idx, row_data in enumerate(comparison_data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(0xf5, 0xf5, 0xf5)

# 第9页：趋势预测
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "2025-2027年市场趋势预测", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)
add_title_shape(slide, "租金与空置率预测走势", Inches(0.5), Inches(1.35), Inches(6), Inches(0.4), font_size=15, color=PRIMARY_COLOR)

forecast_data = {'平均租金': [95, 90, 92, 96], '空置率(×4)': [94, 96, 84, 72]}
chart_img = create_chart_image(forecast_data, 'line', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.8), width=Inches(6.2))
add_body_text(slide, "注：空置率数值已×4缩放以便同图对比", Inches(0.5), Inches(5.7), Inches(5), Inches(0.3), font_size=10, color=GRAY_TEXT)

add_title_shape(slide, "分年度预测数据", Inches(7), Inches(1.35), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
rows, cols = 5, 4
table = slide.shapes.add_table(rows, cols, Inches(6.8), Inches(1.8), Inches(6), Inches(3.2)).table
headers = ['年份', '新增供应(万㎡)', '平均租金(元/㎡)', '空置率']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = SECONDARY_COLOR
    p = cell.text_frame.paragraphs[0]
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

forecast_table = [['2024', '38', '95', '23.5%'], ['2025E', '35', '90', '24.0%'], ['2026E', '28', '92', '21.0%'], ['2027E', '22', '96', '18.0%']]
for row_idx, row_data in enumerate(forecast_table, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER

add_title_shape(slide, "关键预测结论", Inches(7), Inches(5.1), Inches(5.5), Inches(0.4), font_size=15, color=PRIMARY_COLOR)
predictions = ["2025年：空置率预计见顶(24%)", "2026年：市场进入供需平衡期", "2027年：回归健康发展轨道"]
add_bullet_points(slide, predictions, Inches(6.8), Inches(5.5), Inches(6), Inches(1.6), font_size=13)

# 第10页：投资建议
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_COLOR)
add_title_shape(slide, "投资策略建议", Inches(0.6), Inches(0.22), Inches(10), Inches(0.7), font_size=32, color=WHITE)

sections = [("业主/开发商", ["短期：灵活定价，延长免租期至3-6个月", "中期：提升物业服务，考虑楼宇升级", "长期：打造品牌，引入产业链企业"], PRIMARY_COLOR), ("租户企业", ["当前是较好谈判窗口期", "可争取5-10%租金优惠", "建议签3-5年长租锁定低价", "关注天府新区政策优惠"], SECONDARY_COLOR), ("投资者", ["优选核心地段优质资产", "关注2026年后市场回暖", "考虑困境资产投资机会", "重视长期价值而非短期"], ACCENT_COLOR)]
for i, (title, items, color) in enumerate(sections):
    left = Inches(0.4 + i * 4.3)
    add_decorative_bar(slide, left, Inches(1.35), Inches(4.1), Inches(0.55), color)
    add_title_shape(slide, title, left + Inches(0.15), Inches(1.4), Inches(3.8), Inches(0.45), font_size=18, color=WHITE)
    add_bullet_points(slide, items, left + Inches(0.15), Inches(2.05), Inches(3.9), Inches(4.5), font_size=13)

# 第11页：总结
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, PRIMARY_COLOR)
add_decorative_bar(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(0.08), ACCENT_COLOR)
add_title_shape(slide, "核心观点总结", Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), font_size=40, color=WHITE)

key_points = ["市场进入调整期，供应高峰已过，新增供应逐年递减", "区域分化明显，核心区域表现稳健，新兴区域面临培育期", "租户市场特征明显，企业有更多选择和议价空间", "金融、科技、专业服务仍是需求三大主力", "2026年后市场有望回归健康发展，长期前景乐观"]
for i, point in enumerate(key_points):
    num = f"0{i+1}"
    top = Inches(2.3 + i * 0.9)
    add_title_shape(slide, num, Inches(0.8), top, Inches(0.8), Inches(0.6), font_size=28, color=ACCENT_COLOR)
    add_body_text(slide, point, Inches(1.7), top + Inches(0.1), Inches(10.5), Inches(0.7), font_size=17, color=WHITE)

# 第12页：感谢
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, PRIMARY_COLOR)
add_decorative_bar(slide, Inches(4), Inches(2.8), Inches(5.333), Inches(0.08), ACCENT_COLOR)
add_title_shape(slide, "感谢聆听", Inches(0), Inches(3.2), Inches(13.333), Inches(1), font_size=54, color=WHITE, align=PP_ALIGN.CENTER)
add_title_shape(slide, "欢迎交流讨论", Inches(0), Inches(4.2), Inches(13.333), Inches(0.8), font_size=28, bold=False, color=RgbColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER)

# 保存
output_path = os.path.join(output_dir, "成都写字楼市场分析报告_数据增强版.pptx")
prs.save(output_path)
print(f"PPT已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
