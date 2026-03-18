#!/usr/bin/env python3
"""
Generate well-formatted PPT with structured content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets):
    """bullets is list of (level, text) tuples"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content with proper bullet hierarchy
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, (level, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.level = level
        if level == 0:  # Main bullet
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.space_before = Pt(12)
        elif level == 1:  # Sub bullet
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(4)
        else:  # Detail bullet
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_before = Pt(2)
        p.space_after = Pt(6)
    return slide

# ===== GENERATE SLIDES =====

# Slide 1: Title
add_title_slide(prs, "皇庭国际中心数字化营销服务建议书", "打造CBD商务新地标 · 数字化营销全案解决方案")

# Slide 2: Contents
add_content_slide(prs, "目录", [
    (0, "一、项目概况与定位"),
    (1, "项目基本情况、核心优势、目标客群分析"),
    (0, "二、市场环境分析"),
    (1, "CBD商务环境、竞争格局、市场机会洞察"),
    (0, "三、数字化营销策略"),
    (1, "官网小程序、短视频直播、精准广告投放、私域运营"),
    (0, "四、第一太平戴维斯渠道策略"),
    (1, "全球资源优势、渠道合作体系、激励机制"),
    (0, "五、推广执行计划"),
    (1, "三阶段执行方案、时间节点、预期成果"),
    (0, "六、预期效果与合作方案"),
    (1, "量化目标、费用结构、双方职责分工"),
])

# Slide 3: Project Overview
add_content_slide(prs, "项目概况——皇庭国际中心项目介绍", [
    (0, "【项目基本情况】"),
    (1, "位于城市CBD核心区，皇庭集团旗舰商务综合体"),
    (1, "总建筑面积12万㎡：甲级写字楼8万㎡ + 商业配套2万㎡ + 停车场2万㎡"),
    (1, "地上38层，地下3层，预计2024年第三季度交付"),
    (1, "定位为5A甲级智能商务办公空间，LEED金级预认证"),
    (0, "【建筑品质优势】"),
    (1, "国际知名设计事务所操刀，全玻璃幕墙+LOW-E节能设计"),
    (1, "单层面积1500-2000㎡，层高4.2米（市场平均3.8米）"),
    (1, "智能化配置：人脸识别门禁、智能停车、5G全覆盖、楼宇自控系统"),
    (0, "【区位交通优势】"),
    (1, "双地铁交汇上盖：距1号线、3号线交汇站仅200米，步行3分钟"),
    (1, "30分钟可达国际机场，15分钟可达高铁站"),
    (1, "周边配套：五星级酒店8家、购物中心6个、米其林餐厅12家"),
    (0, "【物业服务】"),
    (1, "第一太平戴维斯提供24小时管家式国际化物业服务"),
])

# Slide 4: Target Customers
add_content_slide(prs, "目标客群定位——精准锁定高价值客户", [
    (0, "【三大目标客群】"),
    (1, "金融机构（占比35%）"),
    (2, "银行分行、证券公司营业部、保险公司区域中心、基金公司"),
    (2, "特点：对办公环境形象要求高，注重楼宇品质，租金支付能力强"),
    (1, "科技企业（占比25%）"),
    (2, "互联网公司、软件/AI/大数据企业总部或区域中心"),
    (2, "特点：员工年轻，注重智能化程度和办公舒适度"),
    (1, "专业服务机构（占比20%）"),
    (2, "律师事务所、会计师事务所、咨询公司、广告公司"),
    (2, "特点：重视商务氛围和周边配套，商务接待需求频繁"),
    (1, "其他企业（占比20%）：贸易公司、制造业总部等"),
    (0, "【客户典型特征】"),
    (1, "年营业收入5000万元以上，员工规模100-500人"),
    (1, "日租金预算8-12元/㎡，选址决策周期3-6个月"),
    (1, "决策人：CEO、COO或行政总监"),
])

# Slide 5: Market Analysis
add_content_slide(prs, "市场环境分析——CBD商务区深度洞察", [
    (0, "【区域商务环境】"),
    (1, "城市第一CBD核心区，占地8平方公里"),
    (1, "经济贡献：全市25%的GDP，30%的税收收入"),
    (1, "产业聚集：60%金融机构总部 + 40%科技公司总部 + 50%专业服务机构"),
    (0, "【市场需求分析】"),
    (1, "现有甲级写字楼200万㎡，入驻企业5000+家，就业人口30万"),
    (1, "约30%现有企业有升级办公环境需求，预计150家企业将搬迁/扩租"),
    (1, "每年新增注册企业2000家，约10%符合甲级写字楼标准"),
    (0, "【竞争格局与机会】"),
    (1, "区域内甲级写字楼项目15个，总体空置率8%（供不应求）"),
    (1, "未来12个月内无新增供应，皇庭国际中心抢占市场窗口期"),
    (1, "区域内金融、科技、专业服务三大产业集聚效应显著"),
])

# Slide 6: Competition
add_content_slide(prs, "竞争项目对比分析", [
    (0, "【竞争项目A】位于项目东侧800米，2019年入市"),
    (1, "基本情况：10万㎡，出租率85%，日租金9-11元/㎡"),
    (1, "优势：开发商品牌知名度高，已入驻多家知名企业"),
    (1, "劣势：楼龄5年设施老化，装修风格传统，智能化程度低"),
    (0, "【竞争项目B】位于项目西侧1.2公里，2021年入市"),
    (1, "基本情况：8万㎡，出租率78%，日租金8-10元/㎡"),
    (1, "优势：租金定价较低"),
    (1, "劣势：距地铁站远（步行15分钟），周边配套不完善"),
    (0, "【皇庭国际中心核心优势】"),
    (1, "交通便利性最佳：双地铁交汇上盖，距地铁站仅200米"),
    (1, "建筑品质最新：2024年全新交付，智能化程度最高"),
    (1, "物业服务最优：第一太平戴维斯国际化标准服务"),
    (1, "目标租金10-12元/㎡，综合性价比高"),
])

# Slide 7: Digital Marketing Section
add_section_slide(prs, "数字化营销策略", "全渠道数字化营销体系，精准触达目标客户")

# Slide 8: Website
add_content_slide(prs, "官方网站与小程序生态建设", [
    (0, "【网站建设目标】"),
    (1, "打造高品质响应式官网，适配PC/平板/手机多端"),
    (1, "采用最新前端技术，页面加载速度控制在3秒以内"),
    (0, "【六大核心功能模块】"),
    (1, "VR全景看房系统"),
    (2, "360度全景拍摄大堂、标准层、公共区域"),
    (2, "用户可自由视角查看，身临其境感受项目品质"),
    (1, "楼层平面展示功能"),
    (2, "查看各楼层详细平面图，了解户型分布和可租赁单元"),
    (1, "实时房源查询系统"),
    (2, "实时显示可租面积、租金价格、朝向景观"),
    (2, "支持按面积、价格、楼层筛选"),
    (1, "在线预约看房功能"),
    (2, "填写信息、选择看房时间，自动发送确认短信和邮件"),
    (1, "智能客服机器人（ChatGPT技术）"),
    (2, "7×24小时在线解答常见问题，自动转接人工客服"),
    (1, "数据分析系统"),
    (2, "埋点追踪用户行为，分析关注热点，优化营销策略"),
])

# Save
prs.save('/Users/bob/.openclaw/workspace/皇庭国际中心_排版优化版.pptx')
print("Done! 8 slides with structured formatting.")
