#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
成都写字楼市场分析报告PPT生成器
包含图表、自动排版、专业设计
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

import io
import os

# 创建输出目录
output_dir = "/Users/bob/.openclaw/workspace"
os.makedirs(output_dir, exist_ok=True)

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
PRIMARY_COLOR = RgbColor(0x1a, 0x5f, 0x9e)  # 深蓝
SECONDARY_COLOR = RgbColor(0x2e, 0x8b, 0x57)  # 深绿
ACCENT_COLOR = RgbColor(0xe6, 0x8a, 0x00)  # 橙色
LIGHT_BG = RgbColor(0xf5, 0xf5, 0xf5)  # 浅灰背景
DARK_TEXT = RgbColor(0x33, 0x33, 0x33)  # 深灰文字
WHITE = RgbColor(0xff, 0xff, 0xff)

def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=WHITE):
    """添加标题文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_body_text(slide, text, left, top, width, height, font_size=18, color=DARK_TEXT):
    """添加正文文本"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_bullet_points(slide, items, left, top, width, height, font_size=16):
    """添加项目符号列表"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)
    return shape

def create_chart_image(data, chart_type='bar', title='', colors=None):
    """创建图表图片"""
    fig, ax = plt.subplots(figsize=(6, 4), dpi=150)
    
    if colors is None:
        colors = ['#1a5f9e', '#2e8b57', '#e68a00', '#c44e52', '#8c564b']
    
    if chart_type == 'bar':
        categories = list(data.keys())
        values = list(data.values())
        bars = ax.bar(categories, values, color=colors[:len(categories)])
        ax.set_ylabel('数值', fontsize=11)
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}',
                   ha='center', va='bottom', fontsize=9)
    
    elif chart_type == 'line':
        for label, values in data.items():
            years = list(range(2021, 2026))
            ax.plot(years[:len(values)], values, marker='o', linewidth=2, label=label)
        ax.set_ylabel('数值', fontsize=11)
        ax.legend()
        ax.grid(True, alpha=0.3)
    
    elif chart_type == 'pie':
        sizes = list(data.values())
        labels = list(data.keys())
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors[:len(sizes)],
               startangle=90)
        ax.axis('equal')
    
    elif chart_type == 'horizontal_bar':
        categories = list(data.keys())
        values = list(data.values())
        bars = ax.barh(categories, values, color=colors[:len(categories)])
        ax.set_xlabel('数值', fontsize=11)
        for i, bar in enumerate(bars):
            width = bar.get_width()
            ax.text(width, bar.get_y() + bar.get_height()/2.,
                   f'{width:.1f}',
                   ha='left', va='center', fontsize=9)
    
    ax.set_title(title, fontsize=13, fontweight='bold', pad=15)
    plt.tight_layout()
    
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight', facecolor='white')
    img_stream.seek(0)
    plt.close()
    return img_stream

def add_background_shape(slide, color):
    """添加背景色块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # 移到最底层
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_decorative_bar(slide, left, top, width, height, color):
    """添加装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ============ 第1页：封面 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background_shape(slide, PRIMARY_COLOR)

# 装饰条
add_decorative_bar(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.1), ACCENT_COLOR)

# 主标题
add_title_shape(slide, "成都写字楼市场", Inches(0.8), Inches(2.8), Inches(11), Inches(1), 
                font_size=54, color=WHITE)
add_title_shape(slide, "深度分析报告", Inches(0.8), Inches(3.6), Inches(11), Inches(0.8), 
                font_size=48, color=WHITE)

# 副标题
add_title_shape(slide, "2024-2025年度市场洞察与投资指南", Inches(0.8), Inches(4.6), 
                Inches(11), Inches(0.6), font_size=24, bold=False, color=RgbColor(0xcc, 0xcc, 0xcc))

# 日期
add_title_shape(slide, "2025年3月", Inches(0.8), Inches(6.2), Inches(4), Inches(0.5), 
                font_size=18, bold=False, color=RgbColor(0xaa, 0xaa, 0xaa))

# ============ 第2页：目录 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)

# 标题栏
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "报告目录", Inches(0.6), Inches(0.25), Inches(4), Inches(0.8), 
                font_size=36, color=WHITE)

# 目录内容
sections = [
    ("01", "宏观环境", "城市经济与政策分析"),
    ("02", "市场供给", "存量与增量分析"),
    ("03", "市场需求", "行业结构与租户画像"),
    ("04", "市场表现", "租金与空置率走势"),
    ("05", "区域对比", "四大商务区分析"),
    ("06", "趋势预测", "2025-2027展望"),
    ("07", "投资建议", "策略与风险提示")
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.6 + row * 1.6)
    
    # 编号
    add_title_shape(slide, num, left, top, Inches(1), Inches(0.8), 
                    font_size=32, color=PRIMARY_COLOR)
    # 标题
    add_title_shape(slide, title, left + Inches(1), top, Inches(3), Inches(0.6), 
                    font_size=24, color=DARK_TEXT)
    # 描述
    add_body_text(slide, desc, left + Inches(1), top + Inches(0.5), Inches(4.5), Inches(0.5), 
                  font_size=14, color=RgbColor(0x88, 0x88, 0x88))

# ============ 第3页：城市经济概况 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "城市经济与产业背景", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 左侧关键数据
add_title_shape(slide, "2024年核心经济数据", Inches(0.6), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=20, color=PRIMARY_COLOR)

econ_data = [
    "GDP总量：2.35万亿元",
    "GDP增速：6.0%",
    "常住人口：2140万人",
    "第三产业占比：65.2%",
    "世界500强落户：312家"
]
add_bullet_points(slide, econ_data, Inches(0.6), Inches(2.1), Inches(5), Inches(3), font_size=16)

# 右侧产业结构饼图
add_title_shape(slide, "产业结构分布", Inches(7), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=20, color=PRIMARY_COLOR)

industry_data = {
    '电子信息': 28,
    '装备制造': 22,
    '金融服务': 15,
    '生物医药': 12,
    '其他服务业': 23
}
chart_img = create_chart_image(industry_data, 'pie', '', 
                               ['#1a5f9e', '#2e8b57', '#e68a00', '#c44e52', '#8c564b'])
slide.shapes.add_picture(chart_img, Inches(6.5), Inches(2), width=Inches(6))

# ============ 第4页：市场存量规模 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "市场存量与供应分析", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 存量数据
add_title_shape(slide, "截至2024年底市场存量", Inches(0.6), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=18, color=PRIMARY_COLOR)

stock_data = [
    "甲级写字楼：520万㎡",
    "乙级写字楼：380万㎡",
    "总存量：900万㎡"
]
add_bullet_points(slide, stock_data, Inches(0.6), Inches(2), Inches(4), Inches(1.5), font_size=16)

# 新增供应趋势图
add_title_shape(slide, "历年新增供应趋势", Inches(0.6), Inches(3.5), Inches(5), Inches(0.6), 
                font_size=18, color=PRIMARY_COLOR)

supply_data = {
    '新增供应(万㎡)': [58, 52, 45, 38, 35]
}
chart_img = create_chart_image(supply_data, 'line', '')
slide.shapes.add_picture(chart_img, Inches(0.5), Inches(4), width=Inches(5.5))

# 右侧区域分布
add_title_shape(slide, "存量区域分布", Inches(7), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=18, color=PRIMARY_COLOR)

region_data = {
    '金融城': 120,
    '天府新区': 110,
    '大源': 95,
    '天府广场': 85,
    '其他': 110
}
chart_img = create_chart_image(region_data, 'horizontal_bar', '')
slide.shapes.add_picture(chart_img, Inches(6.5), Inches(2), width=Inches(6))

# ============ 第5页：需求结构 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "租赁需求结构分析", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 左侧饼图
demand_data = {
    '金融业': 35,
    '科技互联网': 28,
    '专业服务': 20,
    '房地产建筑': 10,
    '其他': 7
}
chart_img = create_chart_image(demand_data, 'pie', '需求占比分布')
slide.shapes.add_picture(chart_img, Inches(0.5), Inches(1.5), width=Inches(6))

# 右侧详细说明
add_title_shape(slide, "各行业特征", Inches(7), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=18, color=PRIMARY_COLOR)

industry_details = [
    "金融业(35%)：银行、保险、证券、金融科技",
    "科技互联网(28%)：互联网大厂、本土科技企业",
    "专业服务(20%)：律所、会计所、咨询公司",
    "房地产建筑(10%)：开发商、设计院",
    "其他行业(7%)：制造业、贸易、协会"
]
add_bullet_points(slide, industry_details, Inches(7), Inches(2.1), Inches(5.5), Inches(4.5), font_size=14)

# ============ 第6页：租金与空置率 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "租金水平与空置率", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 分区租金柱状图
add_title_shape(slide, "2024年Q4分区平均租金(元/㎡/月)", Inches(0.6), Inches(1.5), Inches(6), Inches(0.6), 
                font_size=16, color=PRIMARY_COLOR)

rent_data = {
    '金融城': 125,
    '天府广场': 110,
    '大源': 85,
    '天府新区': 75,
    '其他': 70
}
chart_img = create_chart_image(rent_data, 'bar', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(2), width=Inches(6))

# 空置率数据
add_title_shape(slide, "分区空置率对比", Inches(7), Inches(1.5), Inches(6), Inches(0.6), 
                font_size=16, color=PRIMARY_COLOR)

vacancy_data = {
    '天府广场': 16,
    '金融城': 18,
    '大源': 26,
    '其他': 30,
    '天府新区': 32
}
chart_img = create_chart_image(vacancy_data, 'bar', '', ['#c44e52', '#e68a00', '#8c564b', '#2e8b57', '#1a5f9e'])
slide.shapes.add_picture(chart_img, Inches(6.8), Inches(2), width=Inches(6))

# 底部关键指标
add_decorative_bar(slide, Inches(0), Inches(6.2), Inches(13.333), Inches(1.3), LIGHT_BG)
metrics = [
    ("全市平均租金", "95元/㎡/月", "同比-4%"),
    ("甲级空置率", "23.5%", "同比+1.2pp"),
    ("乙级空置率", "28.0%", "同比+0.8pp")
]
for i, (label, value, change) in enumerate(metrics):
    left = Inches(0.8 + i * 4)
    add_title_shape(slide, label, left, Inches(6.3), Inches(3), Inches(0.4), 
                    font_size=14, color=RgbColor(0x66, 0x66, 0x66))
    add_title_shape(slide, value, left, Inches(6.6), Inches(3), Inches(0.5), 
                    font_size=24, color=PRIMARY_COLOR)
    add_body_text(slide, change, left, Inches(7.05), Inches(2), Inches(0.4), 
                  font_size=12, color=RgbColor(0x99, 0x99, 0x99))

# ============ 第7页：区域对比 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "四大商务区对比", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 创建对比表格
# from pptx.enum.table import MSO_ANCHOR as TableAnchor
rows, cols = 5, 6
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5)).table

# 表头
headers = ['区域', '定位', '成熟度', '租金(元/㎡/月)', '空置率', '去化周期']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.color.rgb = WHITE
    paragraph.font.bold = True
    paragraph.font.size = Pt(12)
    paragraph.alignment = PP_ALIGN.CENTER

# 数据
data_rows = [
    ['天府广场', '传统CBD', '成熟', '110', '16%', '12个月'],
    ['金融城', '金融总部', '成熟', '125', '18%', '14个月'],
    ['大源', '科技商务', '较成熟', '85', '26%', '22个月'],
    ['天府新区', '新兴CBD', '发展中', '75', '32%', '30个月']
]

for row_idx, row_data in enumerate(data_rows, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_ALIGN.CENTER
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(0xf0, 0xf5, 0xfa)

# ============ 第8页：未来趋势 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "2025-2027年市场预测", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 左侧趋势图
add_title_shape(slide, "租金与空置率预测", Inches(0.6), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=16, color=PRIMARY_COLOR)

forecast_data = {
    '平均租金': [95, 90, 92, 96],
    '空置率(×4)': [23.5*4, 24*4, 21*4, 18*4]  # 缩放以便同图显示
}
chart_img = create_chart_image(forecast_data, 'line', '')
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(2), width=Inches(6))

# 图例说明
add_body_text(slide, "注：空置率数值已×4缩放以便对比", Inches(0.6), Inches(6), 
              Inches(5), Inches(0.4), font_size=10, color=RgbColor(0x99, 0x99, 0x99))

# 右侧关键预测点
add_title_shape(slide, "关键预测", Inches(7), Inches(1.5), Inches(5), Inches(0.6), 
                font_size=16, color=PRIMARY_COLOR)

predictions = [
    "2025年：空置率预计见顶(24%)",
    "2026年：市场进入平衡期",
    "2027年：回归健康发展轨道",
    "供应逐年递减至30万㎡以下",
    "年净吸纳量回升至40-50万㎡"
]
add_bullet_points(slide, predictions, Inches(7), Inches(2), Inches(5.5), Inches(4), font_size=15)

# ============ 第9页：投资建议 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, WHITE)
add_decorative_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
add_title_shape(slide, "投资策略建议", Inches(0.6), Inches(0.25), Inches(8), Inches(0.8), 
                font_size=32, color=WHITE)

# 三列布局
sections = [
    ("业主/开发商", [
        "灵活定价策略",
        "延长免租期至3-6个月",
        "提升物业服务品质",
        "考虑楼宇升级改造"
    ], PRIMARY_COLOR),
    ("租户", [
        "当前是较好谈判窗口",
        "可争取5-10%租金优惠",
        "签3-5年长租锁定低价",
        "关注天府新区政策优惠"
    ], SECONDARY_COLOR),
    ("投资者", [
        "优选核心地段优质资产",
        "关注2026年后回暖机会",
        "考虑困境资产投资机会",
        "重视长期价值而非短期"
    ], ACCENT_COLOR)
]

for i, (title, items, color) in enumerate(sections):
    left = Inches(0.5 + i * 4.2)
    
    # 标题栏
    add_decorative_bar(slide, left, Inches(1.5), Inches(3.8), Inches(0.6), color)
    add_title_shape(slide, title, left + Inches(0.2), Inches(1.55), Inches(3.4), Inches(0.5), 
                    font_size=18, color=WHITE)
    
    # 内容
    add_bullet_points(slide, items, left + Inches(0.2), Inches(2.3), Inches(3.6), Inches(4), font_size=14)

# ============ 第10页：总结 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, PRIMARY_COLOR)

# 装饰条
add_decorative_bar(slide, Inches(0), Inches(2.2), Inches(13.333), Inches(0.08), ACCENT_COLOR)

# 标题
add_title_shape(slide, "核心观点总结", Inches(0.8), Inches(1.2), Inches(11), Inches(0.8), 
                font_size=40, color=WHITE)

# 核心观点
key_points = [
    "市场进入调整期，供应高峰已过，新增供应逐年递减",
    "区域分化明显，核心区域表现稳健，新兴区域面临培育期",
    "租户市场特征明显，企业有更多选择和议价空间",
    "金融、科技、专业服务仍是需求三大主力",
    "2026年后市场有望回归健康发展，长期前景乐观"
]

for i, point in enumerate(key_points):
    num = f"0{i+1}"
    top = Inches(2.5 + i * 0.9)
    
    # 编号
    add_title_shape(slide, num, Inches(0.8), top, Inches(0.8), Inches(0.6), 
                    font_size=28, color=ACCENT_COLOR)
    # 内容
    add_body_text(slide, point, Inches(1.7), top + Inches(0.1), Inches(10.5), Inches(0.7), 
                  font_size=17, color=WHITE)

# ============ 第11页：感谢 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide, PRIMARY_COLOR)

# 装饰条
add_decorative_bar(slide, Inches(4), Inches(2.8), Inches(5.333), Inches(0.08), ACCENT_COLOR)

# 感谢文字
add_title_shape(slide, "感谢聆听", Inches(0), Inches(3.2), Inches(13.333), Inches(1), 
                font_size=54, color=WHITE)
add_title_shape(slide, "欢迎交流讨论", Inches(0), Inches(4.2), Inches(13.333), Inches(0.8), 
                font_size=28, bold=False, color=RgbColor(0xcc, 0xcc, 0xcc))

# 保存文件
output_path = os.path.join(output_dir, "成都写字楼市场分析报告.pptx")
prs.save(output_path)
print(f"PPT已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
