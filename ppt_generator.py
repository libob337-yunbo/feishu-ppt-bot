# PPT Generator Module
# 使用 python-pptx 生成商务简约风格 PPT

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData, CategoryChartData
import os
import json
import re
from datetime import datetime
import io
import base64
try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')  # 非交互式后端
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

# 商务简约配色方案
COLORS = {
    'primary': RGBColor(0x1E, 0x3A, 0x5F),      # 深蓝 #1E3A5F
    'secondary': RGBColor(0x2E, 0x5C, 0x8A),    # 中蓝 #2E5C8A
    'accent': RGBColor(0x4A, 0x90, 0xA4),       # 青蓝 #4A90A4
    'text_dark': RGBColor(0x33, 0x33, 0x33),    # 深灰 #333333
    'text_light': RGBColor(0x66, 0x66, 0x66),   # 浅灰 #666666
    'white': RGBColor(0xFF, 0xFF, 0xFF),        # 白色
    'bg_light': RGBColor(0xF5, 0xF7, 0xFA),     # 浅灰背景
}

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # 设置幻灯片尺寸为 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self, title, subtitle=""):
        """添加封面页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            Inches(13.333), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary']
        shape.line.fill.background()
        
        # 添加装饰线条
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.5), 
            Inches(2), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS['accent']
        line.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        
        # 添加副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['accent']
            p.font.name = "Microsoft YaHei"
            
        # 添加日期
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(4), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        
    def add_content_slide(self, title, content_list):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(0.15)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Microsoft YaHei"
        
        # 添加内容
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理列表项
            if item.startswith('•') or item.startswith('-') or item.startswith('*'):
                p.text = "    " + item[1:].strip()
                p.level = 0
            elif re.match(r'^\d+[\.\)]\s', item):
                p.text = item
                p.level = 0
            else:
                p.text = "    • " + item
                p.level = 0
                
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_dark']
            p.font.name = "Microsoft YaHei"
            p.space_before = Pt(12)
            p.space_after = Pt(6)
            
    def add_section_slide(self, section_title):
        """添加章节分隔页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['secondary']
        bg.line.fill.background()
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
    def add_chart_slide(self, title, chart_data=None, chart_type='bar', categories=None, values=None, series_name='数据'):
        """添加图表页（支持多种图表类型）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(0.15)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Microsoft YaHei"
        
        # 使用原生PPT图表
        chart_values = values if values is not None else chart_data
        if categories and chart_values:
            self._add_native_chart(slide, chart_type, categories, chart_values, series_name)
        else:
            # 回退到文本展示
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, data in enumerate(chart_data or []):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{data}"
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['text_dark']
                p.font.name = "Microsoft YaHei"
                p.space_before = Pt(8)
    
    def _add_native_chart(self, slide, chart_type, categories, values, series_name):
        """添加PPT原生图表"""
        # 图表类型映射
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'pie': XL_CHART_TYPE.PIE,
            'line': XL_CHART_TYPE.LINE,
            'doughnut': XL_CHART_TYPE.DOUGHNUT,
        }
        
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # 准备图表数据
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(series_name, values)
        
        # 添加图表到幻灯片
        x, y, cx, cy = Inches(1), Inches(1.8), Inches(11), Inches(5)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data
        ).chart
        
        # 设置图表样式
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # 设置字体
        chart.font.size = Pt(12)
        chart.font.name = "Microsoft YaHei"
    
    def add_kpi_slide(self, title, kpis, layout='grid'):
        """添加KPI展示页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(0.15)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Microsoft YaHei"
        
        if layout == 'large' and kpis:
            # 单个大KPI布局
            kpi = kpis[0]
            
            # 大数字
            value_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{kpi['value']}{kpi.get('unit', '')}"
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
            
            # KPI名称
            name_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.8))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = kpi['name']
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['text_dark']
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
            
            # 变化值
            if 'change' in kpi:
                change_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11), Inches(0.6))
                tf = change_box.text_frame
                p = tf.paragraphs[0]
                p.text = kpi['change']
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['accent'] if '+' in str(kpi['change']) else COLORS['secondary']
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
        else:
            # 2x2网格布局
            positions = [
                (Inches(0.8), Inches(1.5)),
                (Inches(6.8), Inches(1.5)),
                (Inches(0.8), Inches(4.2)),
                (Inches(6.8), Inches(4.2)),
            ]
            
            for i, kpi in enumerate(kpis[:4]):
                if i >= len(positions):
                    break
                    
                x, y = positions[i]
                
                # KPI卡片背景
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                    Inches(5.5), Inches(2.5)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = COLORS['bg_light']
                card.line.fill.background()
                
                # KPI值
                value_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.4), Inches(5), Inches(1))
                tf = value_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"{kpi['value']}{kpi.get('unit', '')}"
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = COLORS['primary']
                p.font.name = "Microsoft YaHei"
                
                # KPI名称
                name_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.4), Inches(5), Inches(0.5))
                tf = name_box.text_frame
                p = tf.paragraphs[0]
                p.text = kpi['name']
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text_dark']
                p.font.name = "Microsoft YaHei"
                
                # 变化值
                if 'change' in kpi:
                    change_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.9), Inches(5), Inches(0.4))
                    tf = change_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = kpi['change']
                    p.font.size = Pt(12)
                    p.font.color.rgb = COLORS['accent'] if '+' in str(kpi['change']) else COLORS['secondary']
                    p.font.name = "Microsoft YaHei"
    
    def add_matplotlib_chart_slide(self, title, matplotlib_fig):
        """添加matplotlib图表页（将matplotlib图表转为图片插入）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(0.15)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "Microsoft YaHei"
        
        # 将matplotlib图表保存为图片并插入
        if matplotlib_fig:
            img_stream = io.BytesIO()
            matplotlib_fig.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
            img_stream.seek(0)
            
            slide.shapes.add_picture(
                img_stream, 
                Inches(1), Inches(1.5), 
                width=Inches(11)
            )
            plt.close(matplotlib_fig)
    
    def create_pie_chart(self, labels, values, title="", colors=None):
        """创建matplotlib饼图并返回figure对象"""
        if not MATPLOTLIB_AVAILABLE:
            return None
            
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # 使用配色方案
        pie_colors = colors or ['#1E3A5F', '#2E5C8A', '#4A90A4', '#7FB3D5', '#B8D4E3']
        
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, autopct='%1.1f%%',
            colors=pie_colors[:len(values)],
            startangle=90, textprops={'fontsize': 12}
        )
        
        # 设置标题
        if title:
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        
        plt.tight_layout()
        return fig
    
    def create_bar_chart(self, categories, values, title="", horizontal=False):
        """创建matplotlib柱状图并返回figure对象"""
        if not MATPLOTLIB_AVAILABLE:
            return None
            
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # 使用主色调
        bar_color = '#1E3A5F'
        
        if horizontal:
            bars = ax.barh(categories, values, color=bar_color)
            ax.set_xlabel('数值', fontsize=12)
        else:
            bars = ax.bar(categories, values, color=bar_color)
            ax.set_ylabel('数值', fontsize=12)
            plt.xticks(rotation=15, ha='right')
        
        # 添加数值标签
        for bar in bars:
            if horizontal:
                width = bar.get_width()
                ax.text(width, bar.get_y() + bar.get_height()/2, 
                       f'{width}', ha='left', va='center', fontsize=10)
            else:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2, height,
                       f'{height}', ha='center', va='bottom', fontsize=10)
        
        # 设置标题
        if title:
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        
        # 设置网格
        ax.grid(axis='y', alpha=0.3, linestyle='--')
        ax.set_axisbelow(True)
        
        plt.tight_layout()
        return fig
    
    def create_line_chart(self, x_data, y_data, title="", xlabel="", ylabel=""):
        """创建matplotlib折线图并返回figure对象"""
        if not MATPLOTLIB_AVAILABLE:
            return None
            
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # 使用主色调
        line_color = '#1E3A5F'
        
        ax.plot(x_data, y_data, marker='o', linewidth=2.5, 
                color=line_color, markersize=8, markerfacecolor='#4A90A4')
        
        # 填充区域
        ax.fill_between(x_data, y_data, alpha=0.3, color=line_color)
        
        # 设置标签
        if xlabel:
            ax.set_xlabel(xlabel, fontsize=12)
        if ylabel:
            ax.set_ylabel(ylabel, fontsize=12)
        if title:
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        
        # 设置网格
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.set_axisbelow(True)
        
        plt.tight_layout()
        return fig
                
    def add_end_slide(self, message="谢谢观看"):
        """添加结束页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['primary']
        bg.line.fill.background()
        
        # 结束语
        msg_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = msg_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
    def parse_outline(self, outline_text):
        """解析大纲文本，提取页面结构"""
        pages = []
        current_page = None
        
        lines = outline_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 检测标题（# 或数字开头）
            if line.startswith('#') or re.match(r'^\d+[\.\)]', line):
                if current_page:
                    pages.append(current_page)
                title = re.sub(r'^[#\d\.\)\s]+', '', line)
                current_page = {'title': title, 'content': []}
            elif current_page and (line.startswith('•') or line.startswith('-') or line.startswith('*')):
                current_page['content'].append(line)
            elif current_page and len(line) > 0:
                current_page['content'].append(line)
                
        if current_page:
            pages.append(current_page)
            
        return pages
        
    def generate_from_outline(self, topic, outline_text, detail_text=""):
        """根据大纲生成完整PPT"""
        # 封面
        self.add_title_slide(topic, "AI 智能生成")
        
        # 解析大纲
        pages = self.parse_outline(outline_text)
        
        # 目录页
        if len(pages) > 0:
            self.add_section_slide("目录")
            toc_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            
            # 顶部色条
            header = toc_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(13.333), Inches(0.15)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = COLORS['primary']
            header.line.fill.background()
            
            title_box = toc_slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "目录"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.font.name = "Microsoft YaHei"
            
            content_box = toc_slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, page in enumerate(pages[:6]):  # 最多显示6个章节
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{i+1}. {page['title']}"
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['text_dark']
                p.font.name = "Microsoft YaHei"
                p.space_before = Pt(16)
        
        # 内容页
        for i, page in enumerate(pages):
            if i % 3 == 0 and i > 0:  # 每3页添加一个章节分隔
                self.add_section_slide(page['title'])
            self.add_content_slide(page['title'], page['content'])
            
        # 结束页
        self.add_end_slide()
        
    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        return filename


def generate_ppt_file(topic, outline, detail="", output_dir="/tmp"):
    """生成PPT文件的便捷函数"""
    generator = PPTGenerator()
    generator.generate_from_outline(topic, outline, detail)
    
    filename = f"{output_dir}/PPT_{topic[:20]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    generator.save(filename)
    return filename


if __name__ == "__main__":
    # 测试
    test_outline = """
# 成都写字楼市场分析

## 市场概况
• 成都写字楼总存量超过500万平方米
• 主要分布在高新区、天府新区、锦江区
• 2024年整体空置率约18%

## 租金走势
• 平均租金约80-120元/㎡/月
• 核心商务区租金保持稳定
• 新兴区域租金呈下降趋势

## 投资建议
• 关注TOD项目周边写字楼
• 优选地铁沿线物业
• 考虑产业聚集效应
"""
    
    filename = generate_ppt_file("成都写字楼市场分析", test_outline)
    print(f"PPT已生成: {filename}")
