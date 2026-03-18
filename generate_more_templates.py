#!/usr/bin/env python3
"""
Generate additional PPT templates for ppt-generator skill
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_tech_template():
    """科技风模板 - 适合科技、互联网、创新项目"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 主色调：科技蓝+青色
    primary_color = RGBColor(0, 150, 255)
    accent_color = RGBColor(0, 200, 200)
    dark_bg = RGBColor(15, 20, 35)
    
    # 标题页
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = dark_bg
    bg.line.fill.background()
    
    # 装饰线条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.5), Inches(11.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Presentation Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER
    
    # 内容页
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 248, 250)
    bg.line.fill.background()
    
    # 顶部装饰条
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header.fill.solid()
    header.fill.fore_color.rgb = primary_color
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Slide Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # 内容卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(220, 230, 240)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(5))
    tf = content_box.text_frame
    tf.text = "Content"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 70, 80)
    
    prs.save('~/.openclaw/skills/ppt-generator/assets/templates/tech.pptx')
    print("Created: tech.pptx (科技风)")

def create_elegant_template():
    """优雅风模板 - 适合金融、高端商务、奢侈品"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 主色调：金色+深灰
    gold_color = RGBColor(180, 140, 60)
    dark_gray = RGBColor(50, 50, 55)
    cream_bg = RGBColor(252, 250, 245)
    
    # 标题页
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = dark_gray
    bg.line.fill.background()
    
    # 金色装饰边框
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5))
    border.fill.background()
    border.line.color.rgb = gold_color
    border.line.width = Pt(2)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "Presentation Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = gold_color
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # 内容页
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = cream_bg
    bg.line.fill.background()
    
    # 左侧金色条
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = gold_color
    sidebar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Slide Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = dark_gray
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.text = "Content"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    prs.save('~/.openclaw/skills/ppt-generator/assets/templates/elegant.pptx')
    print("Created: elegant.pptx (优雅风)")

def create_vibrant_template():
    """活力风模板 - 适合创意、营销、年轻品牌"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 主色调：橙红渐变
    orange = RGBColor(255, 120, 50)
    coral = RGBColor(255, 80, 100)
    
    # 标题页
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 渐变背景效果（用形状模拟）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = orange
    bg.line.fill.background()
    
    # 装饰圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = coral
    circle.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 150, 80)
    circle2.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "Presentation Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 230, 220)
    p.alignment = PP_ALIGN.CENTER
    
    # 内容页
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 250, 248)
    bg.line.fill.background()
    
    # 顶部彩色条
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
    header.fill.solid()
    header.fill.fore_color.rgb = orange
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Slide Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = orange
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.text = "Content"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    prs.save('~/.openclaw/skills/ppt-generator/assets/templates/vibrant.pptx')
    print("Created: vibrant.pptx (活力风)")

def create_nature_template():
    """自然风模板 - 适合环保、健康、生活方式"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 主色调：森林绿+米白
    forest_green = RGBColor(60, 120, 80)
    sage_green = RGBColor(140, 170, 130)
    cream = RGBColor(250, 248, 240)
    
    # 标题页
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = forest_green
    bg.line.fill.background()
    
    # 底部装饰
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6), prs.slide_width, Inches(1.5))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = sage_green
    bottom.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Presentation Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(220, 235, 220)
    p.alignment = PP_ALIGN.CENTER
    
    # 内容页
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = cream
    bg.line.fill.background()
    
    # 顶部绿色条
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header.fill.solid()
    header.fill.fore_color.rgb = forest_green
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Slide Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = forest_green
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.text = "Content"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(70, 80, 70)
    
    prs.save('~/.openclaw/skills/ppt-generator/assets/templates/nature.pptx')
    print("Created: nature.pptx (自然风)")

def create_purple_template():
    """紫罗兰模板 - 适合女性、美妆、时尚"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 主色调：紫罗兰+浅紫
    purple = RGBColor(120, 80, 160)
    light_purple = RGBColor(200, 180, 220)
    lavender = RGBColor(245, 240, 250)
    
    # 标题页
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = purple
    bg.line.fill.background()
    
    # 装饰波浪（用椭圆模拟）
    wave = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5.5), Inches(17.333), Inches(4))
    wave.fill.solid()
    wave.fill.fore_color.rgb = light_purple
    wave.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Presentation Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = "Subtitle"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(230, 220, 240)
    p.alignment = PP_ALIGN.CENTER
    
    # 内容页
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = lavender
    bg.line.fill.background()
    
    # 顶部紫色条
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header.fill.solid()
    header.fill.fore_color.rgb = purple
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Slide Title"
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = purple
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.text = "Content"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 70, 90)
    
    prs.save('~/.openclaw/skills/ppt-generator/assets/templates/purple.pptx')
    print("Created: purple.pptx (紫罗兰风)")

# Run
if __name__ == '__main__':
    import os
    os.makedirs('~/.openclaw/skills/ppt-generator/assets/templates', exist_ok=True)
    
    create_tech_template()
    create_elegant_template()
    create_vibrant_template()
    create_nature_template()
    create_purple_template()
    
    print("\nAll templates created successfully!")
    print("Available templates:")
    print("  - business.pptx (商务蓝)")
    print("  - minimal.pptx (极简白)")
    print("  - dark.pptx (深色)")
    print("  - tech.pptx (科技风)")
    print("  - elegant.pptx (优雅风-金)")
    print("  - vibrant.pptx (活力风-橙)")
    print("  - nature.pptx (自然风-绿)")
    print("  - purple.pptx (紫罗兰)")
