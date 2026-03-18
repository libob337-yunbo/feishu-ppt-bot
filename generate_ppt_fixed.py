#!/usr/bin/env python3
"""
Generate detailed PPT with full paragraph content - Fixed version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, subtitle=""):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, content):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    paragraphs = content.split('\n\n')
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para.strip()
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
        p.line_spacing = 1.2
    return slide

# ===== GENERATE SLIDES =====

# Slide 1
add_title_slide(prs, "皇庭国际中心数字化营销服务建议书", "打造CBD商务新地标 · 数字化营销全案解决方案")

# Slide 2 - 目录
add_content_slide(prs, "目录", 
"项目概况与定位——详细介绍皇庭国际中心的项目基本情况、核心优势和目标客群定位分析\n\n市场环境分析——深入分析CBD商务区市场环境、竞争格局和市场需求机会\n\n数字化营销策略——全面阐述官网小程序建设、短视频直播营销、精准广告投放和私域流量运营策略\n\n第一太平戴维斯渠道策略——详细介绍第一太平戴维斯的全球资源优势、渠道合作体系和激励机制\n\n推广执行计划——分阶段说明数字化基建期、全面推广期和精准转化期的具体执行方案\n\n预期效果与目标——明确数字化营销效果目标、招商目标和投资回报目标\n\n合作方案——说明服务内容、费用结构、第一太平戴维斯专属条款和双方职责分工")

# Slide 3 - 项目概况
text3 = """皇庭国际中心位于城市CBD核心区，是皇庭集团倾力打造的旗舰商务综合体项目。项目总建筑面积12万平方米，其中地上38层，地下3层，业态包括甲级写字楼8万平方米、商业配套2万平方米、停车场2万平方米。项目预计2024年第三季度正式交付使用，定位为5A甲级智能商务办公空间，目标客户为金融、科技、专业服务类中大型企业。

在建筑品质方面，项目由国际知名建筑设计事务所操刀设计，采用全玻璃幕墙配合LOW-E节能设计，不仅外观现代大气，更具有出色的节能效果。单层面积1500-2000平方米，层高达到4.2米，远超市场平均3.8米的标准，为企业提供更开阔的办公空间感受。项目已获得LEED金级预认证，配备了智能化楼宇管理系统，包括人脸识别门禁系统、智能停车引导系统、5G网络全覆盖、楼宇自控系统等先进设施。

交通区位是本项目的一大核心优势。项目距离地铁1号线、3号线交汇站仅200米，步行3分钟即可到达，真正实现地铁上盖。通过城市快速路网，30分钟可达国际机场，15分钟可达高铁站。项目周边聚集了城市最高端的商业配套资源，包括五星级酒店8家、大型购物中心6个、米其林星级餐厅12家，完全满足企业商务接待和员工日常生活需求。物业服务由全球领先的第一太平戴维斯提供，这家拥有160年历史的国际物业服务企业将为入驻企业提供24小时管家式专业服务。"""
add_content_slide(prs, "项目概况——皇庭国际中心项目介绍", text3)

# Slide 4 - 目标客群
text4 = """根据项目所在的CBD核心区位和产品高端定位，我们将目标客群精准定位为三大类企业客户。第一类是金融机构，具体包括银行分行、证券公司营业部、保险公司区域中心、基金公司等，这类客户对办公环境的形象展示要求极高，非常注重楼宇品质和商务氛围，具备很强的租金支付能力，在我们的目标客户中占比约35%。

第二类是科技企业，涵盖互联网公司、软件开发企业、人工智能公司、大数据企业等新兴科技行业的总部或区域中心。这类客户的员工群体相对年轻，更加注重办公环境的智能化程度和舒适度，对5G网络、智能办公设施等配置要求较高，在我们的目标客户中占比约25%。

第三类是专业服务机构，包括律师事务所、会计师事务所、管理咨询公司、广告公司、设计公司等。这类客户的业务性质决定了他们对商务氛围和周边配套非常重视，需要频繁进行商务接待和会议活动，在我们的目标客户中占比约20%。此外还有贸易公司、制造业企业总部等其他类型企业占比20%。

这些目标客户的典型特征可以概括为：年营业收入5000万元以上，员工规模100-500人，对办公环境品质要求高且注重企业形象展示，日租金预算在8-12元每平方米之间，选址决策周期通常为3-6个月，选址决策通常由企业CEO、COO或行政总监负责。"""
add_content_slide(prs, "目标客群定位——精准锁定高价值客户群体", text4)

# Slide 5 - 市场分析
text5 = """项目所在的区域是城市第一CBD核心区，规划占地面积约8平方公里，这个区域贡献了全市25%的GDP产出和30%的税收收入，是城市经济发展最活跃、商务活动最密集的核心区域。从产业聚集情况来看，区域内汇集了全市60%的金融机构总部、40%的科技公司总部、50%的专业服务机构，已经形成了金融、科技、专业服务三大主导产业的强大集聚效应，这种产业集聚为企业间的商务合作和资源整合提供了得天独厚的条件。

从企业办公需求角度分析，目前区域内现有甲级写字楼总建筑面积约200万平方米，入驻企业超过5000家，就业人口约30万人，构成了庞大的商务人群。根据我们前期的深入市场调研，约30%的现有企业存在升级办公环境的明确需求，预计有超过150家企业可能在2024-2025年期间考虑搬迁到更优质的办公空间或进行扩租。同时，随着城市产业结构的持续升级优化，每年新增注册企业约2000家，其中约10%符合甲级写字楼的入驻标准，这为项目提供了持续稳定的潜在客源。

从市场竞争格局来看，区域内现有甲级写字楼项目共15个，总体空置率约为8%，处于较低水平，这意味着优质的甲级写字楼项目处于供不应求的状态，是项目入市的良好时机。特别值得注意的是，未来12个月内，区域内暂时没有新的甲级写字楼项目供应计划，这给了皇庭国际中心宝贵的时间窗口来抢占市场份额。"""
add_content_slide(prs, "市场环境分析——CBD商务区深度市场洞察", text5)

# Slide 6 - 竞争分析
text6 = """竞争项目A位于本项目东侧约800米的位置，于2019年入市运营，总建筑面积10万平方米，目前的出租率达到85%，日租金水平在9-11元每平方米。该项目的主要优势在于开发商品牌知名度较高，已经入驻了多家知名企业，形成了一定的集聚效应和品牌背书。但该项目的劣势也很明显，楼龄已经有5年，部分设施开始老化，装修风格相对传统，智能化程度明显不如新项目，难以满足科技型企业对智能化办公的需求。

竞争项目B位于本项目西侧约1.2公里的位置，于2021年入市运营，总建筑面积8万平方米，目前的出租率为78%，日租金水平在8-10元每平方米。该项目的主要优势在于租金定价相对较低，对价格敏感的客户有一定吸引力。但该项目的劣势在于距离地铁站较远，步行需要15分钟，对于依赖公共交通的员工来说通勤便利性较差，同时周边的商业配套也不如核心区完善，对企业的商务接待便利性有一定影响。

相比之下，皇庭国际中心的优势非常明显且难以复制。首先在交通便利性方面，本项目是双地铁交汇上盖，距离地铁站仅200米步行3分钟，是区域内交通最便利的项目。其次在建筑品质方面，本项目是2024年全新交付，建筑设计、设施配置、智能化程度都是最新的。再次在物业服务方面，本项目由第一太平戴维斯提供国际化标准的物业服务，与本地物业有本质区别。综合考虑这些因素，我们将目标租金水平定为日租金10-12元每平方米，虽然略高于竞品，但考虑到产品的综合优势，性价比实际上非常高。"""
add_content_slide(prs, "竞争项目深度对比——凸显皇庭国际中心核心优势", text6)

# Slide 7 - Section
add_section_slide(prs, "数字化营销策略", "全渠道数字化营销体系，精准触达目标客户")

# Slide 8 - Website
text8 = """官方网站是本项目数字化营销体系的核心阵地和流量中枢。我们将投入专业资源打造一个高品质的响应式官方网站，确保网站能够完美适配PC端电脑、平板设备、手机移动端等各种访问终端，为不同场景下的用户提供一致且优质的浏览体验。网站将采用目前最新的前端技术架构进行开发，经过专业优化确保页面平均加载速度控制在3秒以内，这不仅关乎用户体验，也直接影响搜索引擎的排名表现。

网站将配置六大核心功能模块。第一是VR全景看房系统，我们将使用专业级全景相机对项目大堂、标准办公层、公共区域等进行360度全景拍摄，用户可以通过鼠标拖动或手指滑动来自由视角查看，仿佛身临其境般感受项目品质。第二是楼层平面展示功能，用户可以查看各楼层的详细平面图，直观了解不同面积段的户型分布和可租赁单元情况。第三是实时房源查询系统，系统会实时显示各楼层的可租面积、租金价格、朝向景观等信息，并支持用户按面积区间、价格区间、楼层偏好等条件进行灵活筛选。

第四是在线预约看房功能，意向客户可以填写基本信息、选择方便的看房时间，提交后系统会自动发送确认短信和邮件给客户和置业顾问。第五是智能客服机器人，我们将基于ChatGPT技术打造智能客服系统，可以7×24小时在线解答客户关于项目位置、交通配套、租金价格、交付标准等常见问题，当遇到超出机器人处理范围的问题时会自动转接人工客服跟进。第六是数据分析系统，通过在网站部署数据埋点来追踪记录用户的浏览行为轨迹，分析用户关注的热点内容和转化路径，为后续的营销策略优化提供数据支持。"""
add_content_slide(prs, "官方网站与小程序生态建设——打造数字化营销核心阵地", text8)

# Save
prs.save('/Users/bob/.openclaw/workspace/皇庭国际中心_详细内容版.pptx')
print("Done! 8 slides generated.")
